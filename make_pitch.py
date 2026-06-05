from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

GOLD = RGBColor(0xB8, 0x92, 0x4A)
CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x80, 0x80, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Yu Gothic"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s

PH_FILL = RGBColor(0xF2, 0xF2, 0xF2)
PH_LINE = RGBColor(0xD0, 0xD0, 0xD0)

def PH(slide, left_in, top_in, w_in, h_in, label="写真を差し替え"):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))
    box.fill.solid()
    box.fill.fore_color.rgb = PH_FILL
    box.line.color.rgb = PH_LINE
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "［ " + label + " ］"
    r.font.name = FONT
    r.font.size = Pt(14)
    r.font.color.rgb = MUTED
    return box

def T(slide, text, top_in, size, color=CHARCOAL, bold=False, align=PP_ALIGN.CENTER, left_in=0.5, width_in=12.333, italic=False):
    h = (size/72.0) * 1.8 * max(1, text.count("\n")+1)
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

# 1 Title
s = add_slide()
T(s, "LUXERE", 2.3, 130, GOLD, bold=True)
T(s, "出張洗車", 4.6, 22, CHARCOAL)
T(s, "堀本 菖    2026", 6.6, 13, MUTED)

# 2 はじめまして
s = add_slide()
T(s, "はじめまして。", 2.0, 60, CHARCOAL, bold=True, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
T(s, "堀本 菖  （ほりもと しょう）", 3.6, 20, CHARCOAL, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
T(s, "名古屋工業大学  4年  /  21歳", 4.3, 16, MUTED, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
PH(s, 8.2, 1.2, 4.3, 5.1, "顔写真")

# 3 出身
s = add_slide()
T(s, "愛知県  豊川市  出身。", 3.2, 56, CHARCOAL, bold=True)

# 4 通学
s = add_slide()
T(s, "毎朝、", 1.8, 44, MUTED, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
T(s, "片道1時間半。", 2.7, 56, CHARCOAL, bold=True, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
T(s, "実家から、電車で大学へ。", 4.2, 18, MUTED, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
PH(s, 8.2, 1.5, 4.3, 4.5, "電車・通学風景")

# 5 ある日
s = add_slide()
T(s, "ある日、気づいた。", 3.2, 60, CHARCOAL)

# 6 楽しそうな人がいない
s = add_slide()
T(s, "「一人も、\n楽しそうな人がいない。」", 2.2, 52, CHARCOAL, italic=True, bold=True)
T(s, "—— 通勤電車のサラリーマンを見て。", 5.8, 18, MUTED)

# 7 問い
s = add_slide()
T(s, "一度きりの人生を、", 2.4, 48, CHARCOAL)
T(s, "このまま終わらせていいのか。", 3.6, 48, CHARCOAL, bold=True)

# 8 突出したスキルはなかった
s = add_slide()
T(s, "突出したスキルは、", 2.6, 44, MUTED)
T(s, "なかった。", 3.8, 80, CHARCOAL, bold=True)

# 9 でも
s = add_slide()
T(s, "でも、", 1.8, 36, MUTED, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
T(s, "周りの車を、\n洗ってみた。", 2.6, 48, CHARCOAL, bold=True, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
PH(s, 8.2, 1.2, 4.3, 5.1, "洗車中の写真")

# 10 みんな笑ってくれた
s = add_slide()
PH(s, 0.8, 0.8, 11.7, 4.2, "施工後のお客様 / 仕上がり写真")
T(s, "みんな、笑ってくれた。", 5.5, 56, GOLD, bold=True)

# 11 これだ
s = add_slide()
T(s, "「これだ。」", 3.0, 100, CHARCOAL, bold=True)

# 12 一度、諦めた
s = add_slide()
T(s, "大学3年・冬。", 2.0, 20, MUTED)
T(s, "一度、起業した。", 2.8, 48, CHARCOAL)
T(s, "そして、夏。", 4.2, 20, MUTED)
T(s, "—— 一度、諦めた。", 5.0, 36, MUTED, italic=True)

# 13 もう一度
s = add_slide()
T(s, "それでも。", 2.4, 28, MUTED)
T(s, "もう一度。", 3.4, 100, GOLD, bold=True)

# 14 LUXERE
s = add_slide()
T(s, "LUXERE", 1.4, 130, GOLD, bold=True)
T(s, "—— 上質を、あなたのもとへ。", 3.5, 22, CHARCOAL, italic=True)
PH(s, 2.2, 4.5, 8.9, 2.6, "ロゴ / 高級車 / 施工イメージ")

# 15 BNI 感謝
s = add_slide()
T(s, "BNIで、学ばせてもらっています。", 2.8, 40, CHARCOAL, bold=True)
T(s, "大学だけでは、絶対に得られない経験を。", 4.0, 20, MUTED)
T(s, "—— 本当に、ありがとうございます。", 5.4, 22, GOLD, italic=True)

# 16 大切にしていること
s = add_slide()
T(s, "大切にしていること", 0.9, 16, MUTED, align=PP_ALIGN.LEFT, left_in=0.8, width_in=6)
T(s, "丁寧。", 1.8, 54, CHARCOAL, bold=True, align=PP_ALIGN.LEFT, left_in=0.8, width_in=6)
T(s, "誠実。", 3.4, 54, CHARCOAL, bold=True, align=PP_ALIGN.LEFT, left_in=0.8, width_in=6)
T(s, "上質。", 5.0, 54, GOLD, bold=True, align=PP_ALIGN.LEFT, left_in=0.8, width_in=6)
PH(s, 7.5, 1.2, 5.0, 5.1, "ディテール作業の写真")

# 17 課題
s = add_slide()
T(s, "正直、課題もあります。", 2.4, 36, MUTED)
T(s, "認知度。", 3.6, 100, GOLD, bold=True)

# 18 次の一手
s = add_slide()
T(s, "次の一手。", 1.4, 20, MUTED, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
T(s, "ブランディング。", 2.2, 64, CHARCOAL, bold=True, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
T(s, "そして、日本中へ。", 4.0, 28, GOLD, align=PP_ALIGN.LEFT, left_in=0.8, width_in=7)
PH(s, 8.2, 1.5, 4.3, 4.5, "車・日本地図など")

# 19 キャッチコピー
s = add_slide()
T(s, "いつか、こう言われたい。", 2.0, 24, MUTED)
T(s, "「洗車業界と言えば、", 3.2, 60, CHARCOAL, bold=True)
T(s, "堀本。」", 4.4, 100, GOLD, bold=True)

# 20 締め
s = add_slide()
T(s, "ありがとうございました。", 1.6, 48, CHARCOAL, bold=True)
T(s, "一緒に、育てていきませんか。", 3.0, 22, GOLD, italic=True)
PH(s, 4.7, 4.0, 3.9, 2.4, "顔写真 / LUXEREロゴ")
T(s, "堀本 菖    LUXERE    [連絡先]", 6.7, 13, MUTED)

out = r"C:\Users\shang\OneDrive\Desktop\出張洗車\LUXERE_pitch.pptx"
prs.save(out)
print("OK:", out)
